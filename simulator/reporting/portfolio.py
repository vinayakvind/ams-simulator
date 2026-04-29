"""Portfolio reporting helpers for indexed design summaries."""

from __future__ import annotations

from collections.abc import Collection
from datetime import datetime
import json
from pathlib import Path
import tempfile
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from simulator.agents.tech_mapper import TechMapper
from simulator.reporting.design_reference import load_design_reference


PROJECT_ROOT = Path(__file__).resolve().parents[2]


def load_portfolio_entries(
    project_root: Path = PROJECT_ROOT,
    design_ids: Collection[str] | None = None,
) -> list[dict[str, Any]]:
    """Load indexed design manifests and any latest normalized reports."""
    designs_dir = project_root / "designs"
    reports_dir = project_root / "reports"
    selected = {design_id for design_id in (design_ids or [])}
    entries: list[dict[str, Any]] = []

    for manifest_path in sorted(designs_dir.glob("*/design_reference.json")):
        reference = load_design_reference(manifest_path)
        design = reference.get("design", {})
        design_id = str(design.get("id", manifest_path.parent.name))
        if selected and design_id not in selected:
            continue

        report_path = reports_dir / f"{design_id}_regression_latest.json"
        latest_report: dict[str, Any] | None = None
        if report_path.exists():
            latest_report = json.loads(report_path.read_text(encoding="utf-8"))

        tests = latest_report.get("tests", []) if latest_report else []
        grouped: dict[str, list[dict[str, Any]]] = {
            "Analog": [],
            "Digital": [],
            "Top": [],
            "Mixed": [],
            "Other": [],
        }
        for test in tests:
            category = str(test.get("category", "Other"))
            grouped.setdefault(category, []).append(test)

        entries.append(
            {
                "reference": reference,
                "design_id": design_id,
                "design": design,
                "architecture": reference.get("architecture", {}),
                "report": latest_report,
                "tests_by_category": grouped,
            }
        )
    return entries


def build_portfolio_payload(entries: list[dict[str, Any]]) -> dict[str, Any]:
    """Build the normalized JSON payload for the portfolio overview."""
    return {
        "generated_at": datetime.now().isoformat(),
        "supported_technologies": TechMapper.list_technologies(),
        "designs": [
            {
                "id": entry["design_id"],
                "title": entry["design"].get("title", entry["design_id"]),
                "technology": entry["design"].get("technology", ""),
                "summary": entry["design"].get("summary", ""),
                "architecture_summary": entry["architecture"].get("summary", ""),
                "block_count": len(entry["reference"].get("blocks", [])),
                "latest_report": entry["report"].get("summary", {}) if entry.get("report") else {},
                "tests_by_category": entry["tests_by_category"],
            }
            for entry in entries
        ],
    }


def write_portfolio_markdown(entries: list[dict[str, Any]], output_path: Path) -> None:
    """Write the Markdown portfolio summary."""
    lines = [
        "# Design Portfolio Overview",
        "",
        f"Generated: {datetime.now().isoformat()}",
        "",
        "## Supported Technologies",
        "",
    ]
    for tech in TechMapper.list_technologies():
        lines.append(f"- {tech['name']} ({tech['node']}): {tech['description']}")

    lines.extend([
        "",
        "## Design Summary",
        "",
        "| Design | Technology | Blocks | Latest Overall |",
        "|--------|------------|--------|----------------|",
    ])
    for entry in entries:
        summary = entry.get("report", {}).get("summary", {})
        lines.append(
            f"| {entry['design'].get('title', entry['design_id'])} | {entry['design'].get('technology', '')} | "
            f"{len(entry['reference'].get('blocks', []))} | {summary.get('overall', 'NO DATA')} |"
        )

    for entry in entries:
        lines.extend(
            [
                "",
                f"## {entry['design'].get('title', entry['design_id'])}",
                "",
                entry["architecture"].get("summary", "No architecture summary recorded."),
                "",
            ]
        )
        for category in ("Analog", "Mixed", "Digital", "Top"):
            tests = entry["tests_by_category"].get(category, [])
            if not tests:
                continue
            lines.extend(
                [
                    f"### {category} Results",
                    "",
                    "| Case | Block | Status | Details |",
                    "|------|-------|--------|---------|",
                ]
            )
            for test in tests:
                lines.append(
                    f"| {test.get('case_id', '')} | {test.get('block', '')} | {test.get('status', '')} | {test.get('details', '')} |"
                )
            lines.append("")

    output_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _add_title_slide(presentation: Presentation, title: str, subtitle: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(presentation: Presentation, title: str, bullets: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.font.size = Pt(20)


def _add_design_slide(presentation: Presentation, entry: dict[str, Any]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = entry["design"].get("title", entry["design_id"])

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.4), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    summary = [
        f"Technology: {entry['design'].get('technology', 'Not recorded')}",
        f"Blocks: {len(entry['reference'].get('blocks', []))}",
        entry["architecture"].get("summary", "No architecture summary recorded."),
    ]
    for index, line in enumerate(summary):
        paragraph = left_frame.paragraphs[0] if index == 0 else left_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(19 if index < 2 else 17)
        if index == 0:
            paragraph.font.bold = True

    right_box = slide.shapes.add_textbox(Inches(5.1), Inches(1.2), Inches(4.1), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    report_summary = entry.get("report", {}).get("summary", {})
    header = right_frame.paragraphs[0]
    header.text = f"Latest Overall: {report_summary.get('overall', 'NO DATA')}"
    header.font.size = Pt(22)
    header.font.bold = True
    if report_summary.get("overall") == "PASS":
        header.font.color.rgb = RGBColor(46, 125, 50)
    elif report_summary.get("overall"):
        header.font.color.rgb = RGBColor(176, 0, 32)

    for category in ("Analog", "Mixed", "Digital", "Top"):
        tests = entry["tests_by_category"].get(category, [])
        if not tests:
            continue
        cat = right_frame.add_paragraph()
        cat.text = category
        cat.font.size = Pt(18)
        cat.font.bold = True
        for test in tests[:5]:
            paragraph = right_frame.add_paragraph()
            paragraph.text = f"{test.get('case_id', '')}: {test.get('status', '')}"
            paragraph.level = 1
            paragraph.font.size = Pt(16)


def _add_technology_slide(presentation: Presentation) -> None:
    bullets = [
        f"{tech['name']} ({tech['node']}): {tech['description']}" for tech in TechMapper.list_technologies()
    ]
    bullets.append(
        "Imported-source support lets the framework wrap existing SPICE, Verilog, or SystemVerilog designs without regenerating them."
    )
    bullets.append(
        "Technology-aware model cards now come from TechMapper, so the same indexed flow can be retargeted across generic180, generic130, generic65, and bcd180."
    )
    _add_bullet_slide(presentation, "Technology Compatibility", bullets)


def generate_portfolio_artifacts(
    output_ppt: Path | None = None,
    project_root: Path = PROJECT_ROOT,
    design_ids: Collection[str] | None = None,
) -> dict[str, Any]:
    """Generate JSON, Markdown, and PowerPoint portfolio artifacts."""
    entries = load_portfolio_entries(project_root=project_root, design_ids=design_ids)
    reports_dir = project_root / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)

    json_output = reports_dir / "design_portfolio_overview.json"
    markdown_output = reports_dir / "design_portfolio_overview.md"
    ppt_output = Path(output_ppt) if output_ppt else reports_dir / "design_portfolio_overview.pptx"

    payload = build_portfolio_payload(entries)
    json_output.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    write_portfolio_markdown(entries, markdown_output)

    presentation = Presentation()
    _add_title_slide(
        presentation,
        "AMS Design Portfolio Overview",
        f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} | Indexed architecture and run summary",
    )
    _add_technology_slide(presentation)
    _add_bullet_slide(
        presentation,
        "Portfolio Snapshot",
        [
            f"{entry['design'].get('title', entry['design_id'])}: {entry.get('report', {}).get('summary', {}).get('overall', 'NO DATA')}"
            for entry in entries
        ]
        or ["No indexed designs found."],
    )
    for entry in entries:
        _add_design_slide(presentation, entry)

    try:
        presentation.save(str(ppt_output))
    except OSError as exc:
        # Handle potential issues with paths containing spaces or special characters
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
            presentation.save(str(tmp_path))
            tmp_path.replace(ppt_output)
    
    return {
        "json": json_output,
        "markdown": markdown_output,
        "ppt": ppt_output,
        "payload": payload,
    }


__all__ = [
    "build_portfolio_payload",
    "generate_portfolio_artifacts",
    "load_portfolio_entries",
    "write_portfolio_markdown",
]